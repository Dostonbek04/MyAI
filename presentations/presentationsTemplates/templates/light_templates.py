from pptx import Presentation as PPT
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Dizayn uslublari ro‘yxatlari (faqat ma'lumot uchun)
IMAGE_STYLES = [
    "style_focus_left_image", "style_focus_right_image", "style_image_top_text_bottom",
    "style_image_grid_blocks", "style_split_full_image", "style_overlay_image_dark",
    "style_diagonal_layout", "style_focus_image_blend", "style_focus_circle",
    "style_stat_block", "style_side_by_side_card", "style_background_with_quote"
]
TEXT_STYLES = [
    "style_basic_text_centered", "style_bullet_points_left", "style_numbered_steps_right",
    "style_question_vs_answer", "style_vertical_steps", "style_quote_spotlight",
    "style_simple_title_content", "style_two_column_compare", "style_centered_fact",
    "style_timeline_wave", "style_call_to_action_box", "style_title_with_subtext",
    "style_highlight_word", "style_process_chain", "style_stack_blocks",
    "style_minimalist_boxed", "style_definitions"
]

# Light shablonlar ro‘yxati
LIGHT_TEMPLATES = [
    {"name": "Quyosh Nuri", "background_color": "#ffffff", "text_color": "#000000", "style": "style_focus_left_image"},
    {"name": "Ochiq Osmon", "background_color": "#f0f0f0", "text_color": "#333333", "style": "style_image_top_text_bottom"},
    {"name": "Yorug‘ Tong", "background_color": "#f5f5f5", "text_color": "#1a1a1a", "style": "style_basic_text_centered"},
    {"name": "Bahor Shamoli", "background_color": "#e8e8e8", "text_color": "#2b2b2b", "style": "style_split_full_image"},
    {"name": "Sof Oq Dunyo", "background_color": "#fafafa", "text_color": "#222222", "style": "style_bullet_points_left"},
    {"name": "Yorug‘ Ufq", "background_color": "#f2f2f2", "text_color": "#303030", "style": "style_numbered_steps_right"},
    {"name": "Quyoshli Dalalar", "background_color": "#ececec", "text_color": "#252525", "style": "style_question_vs_answer"},
    {"name": "Ochiq Yulduzlar", "background_color": "#ffffff", "text_color": "#1f1f1f", "style": "style_vertical_steps"},
    {"name": "Bahor Yomg‘iri", "background_color": "#e5e5e5", "text_color": "#333333", "style": "style_quote_spotlight"},
    {"name": "Yorug‘ Chiroq", "background_color": "#f7f7f7", "text_color": "#282828", "style": "style_simple_title_content"},
    {"name": "Sof Quyosh", "background_color": "#eeeeee", "text_color": "#2a2a2a", "style": "style_two_column_compare"},
    {"name": "Oq Tuman", "background_color": "#fbfbfb", "text_color": "#1e1e1e", "style": "style_centered_fact"},
    {"name": "Yorug‘ O‘rmon", "background_color": "#e3e3e3", "text_color": "#343434", "style": "style_timeline_wave"},
    {"name": "Quyoshli Tong", "background_color": "#f4f4f4", "text_color": "#272727", "style": "style_call_to_action_box"},
    {"name": "Ochiq Suvlar", "background_color": "#fdfdfd", "text_color": "#202020", "style": "style_title_with_subtext"},
    {"name": "Bahor Nuri", "background_color": "#e9e9e9", "text_color": "#313131", "style": "style_highlight_word"},
    {"name": "Yorug‘ Xazina", "background_color": "#f6f6f6", "text_color": "#292929", "style": "style_process_chain"},
    {"name": "Oq Shahar", "background_color": "#fefefe", "text_color": "#1d1d1d", "style": "style_stack_blocks"},
    {"name": "Quyoshli O‘tlar", "background_color": "#e7e7e7", "text_color": "#323232", "style": "style_minimalist_boxed"},
    {"name": "Yorug‘ Soyalar", "background_color": "#f3f3f3", "text_color": "#2c2c2c", "style": "style_definitions"},
    {"name": "Sof Osmon", "background_color": "#fcfcfc", "text_color": "#212121", "style": "style_focus_right_image"},
    {"name": "Bahor Gullari", "background_color": "#e6e6e6", "text_color": "#333333", "style": "style_image_grid_blocks"},
    {"name": "Yorug‘ Qanotlar", "background_color": "#f5f5f5", "text_color": "#2e2e2e", "style": "style_overlay_image_dark"},
    {"name": "Quyoshli Ufq", "background_color": "#ededed", "text_color": "#262626", "style": "style_diagonal_layout"},
    {"name": "Oq Chiroq", "background_color": "#f9f9f9", "text_color": "#232323", "style": "style_focus_image_blend"},
    {"name": "Yorug‘ Tushlar", "background_color": "#e4e4e4", "text_color": "#343434", "style": "style_focus_circle"},
    {"name": "Sof Bahor", "background_color": "#f8f8f8", "text_color": "#2d2d2d", "style": "style_stat_block"},
    {"name": "Quyoshli Soyalar", "background_color": "#ebebeb", "text_color": "#282828", "style": "style_side_by_side_card"},
    {"name": "Oq Yulduzlar", "background_color": "#fafafa", "text_color": "#202020", "style": "style_background_with_quote"},
    {"name": "Yorug‘ Yomg‘ir", "background_color": "#e2e2e2", "text_color": "#353535", "style": "style_basic_text_centered"},
    {"name": "Quyoshli O‘rmon", "background_color": "#f7f7f7", "text_color": "#2f2f2f", "style": "style_bullet_points_left"},
    {"name": "Sof Ufq", "background_color": "#eeeeee", "text_color": "#262626", "style": "style_numbered_steps_right"},
    {"name": "Bahor Soyasi", "background_color": "#f4f4f4", "text_color": "#313131", "style": "style_question_vs_answer"},
    {"name": "Yorug‘ Suvlar", "background_color": "#fdfdfd", "text_color": "#1f1f1f", "style": "style_vertical_steps"},
    {"name": "Quyoshli Chiroq", "background_color": "#e9e9e9", "text_color": "#323232", "style": "style_quote_spotlight"},
    {"name": "Oq Xotiralar", "background_color": "#f6f6f6", "text_color": "#2a2a2a", "style": "style_simple_title_content"},
    {"name": "Yorug‘ Shahar", "background_color": "#fefefe", "text_color": "#1e1e1e", "style": "style_two_column_compare"},
    {"name": "Sof Yulduz", "background_color": "#e7e7e7", "text_color": "#333333", "style": "style_centered_fact"},
    {"name": "Quyoshli Tong 2", "background_color": "#f5f5f5", "text_color": "#2c2c2c", "style": "style_timeline_wave"},
    {"name": "Bahor Ovozi", "background_color": "#ededed", "text_color": "#272727", "style": "style_call_to_action_box"},
    {"name": "Yorug‘ Ko‘zgu", "background_color": "#f9f9f9", "text_color": "#242424", "style": "style_title_with_subtext"},
    {"name": "Sof Quyosh 2", "background_color": "#e4e4e4", "text_color": "#343434", "style": "style_highlight_word"},
    {"name": "Quyoshli Gullar", "background_color": "#f8f8f8", "text_color": "#2d2d2d", "style": "style_process_chain"},
    {"name": "Oq Soyalar", "background_color": "#ebebeb", "text_color": "#282828", "style": "style_stack_blocks"},
    {"name": "Yorug‘ Shamol", "background_color": "#fafafa", "text_color": "#212121", "style": "style_minimalist_boxed"},
    {"name": "Bahor Yulduzi", "background_color": "#e6e6e6", "text_color": "#323232", "style": "style_definitions"},
    {"name": "Quyoshli Tushlar", "background_color": "#f7f7f7", "text_color": "#2e2e2e", "style": "style_focus_left_image"},
    {"name": "Sof Chiroq", "background_color": "#eeeeee", "text_color": "#262626", "style": "style_image_top_text_bottom"},
    {"name": "Yorug‘ Dalalar", "background_color": "#f4f4f4", "text_color": "#313131", "style": "style_split_full_image"},
    {"name": "Oq Ufqlar", "background_color": "#fdfdfd", "text_color": "#1f1f1f", "style": "style_overlay_image_dark"},
]

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def set_background_color(slide, hex_color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)

def create_light_template(title, slide_titles, slide_texts, image_paths, num_slides, style_index, with_images=False, slide_styles=None):
    if style_index < 0 or style_index >= len(LIGHT_TEMPLATES):
        style_index = 0
    template = LIGHT_TEMPLATES[style_index]
    background_color = template["background_color"]
    text_color = template["text_color"]

    if slide_styles is None or len(slide_styles) != num_slides:
        slide_styles = [template["style"]] * num_slides

    ppt = PPT()

    # Birinchi slayd (sarlavha slaydi)
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    set_background_color(slide, background_color)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.left = Inches(1)
    title_shape.top = Inches(1)
    title_shape.width = Inches(8)
    title_shape.height = Inches(1)
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(text_color)

    # Qolgan slaydlar
    for i in range(num_slides):
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        set_background_color(slide, background_color)

        # Sarlavha har doim yuqorida joylashadi
        title_shape = slide.shapes.title
        title_shape.text = slide_titles[i]
        title_shape.left = Inches(1)
        title_shape.top = Inches(1)
        title_shape.width = Inches(8)
        title_shape.height = Inches(1)
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(text_color)

        has_image = with_images and image_paths[i] is not None
        style = slide_styles[i]

        # Dizayn uslubiga ko‘ra slaydni joylashtirish
        if style == "style_focus_left_image" and has_image:
            # Chapda rasm, o‘ngda matn
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2.5), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(5), Inches(2.5), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_focus_right_image" and has_image:
            # O‘ngda rasm, chapda matn
            slide.shapes.add_picture(image_paths[i], Inches(6), Inches(2.5), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_image_top_text_bottom" and has_image:
            # Yuqorida rasm, pastda matn
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(8), height=Inches(2.5))
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_image_grid_blocks" and has_image:
            # 2x2 rasmli bloklar
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(2), height=Inches(2))
            slide.shapes.add_picture(image_paths[i], Inches(3.5), Inches(2), width=Inches(2), height=Inches(2))
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(4.5), width=Inches(2), height=Inches(2))
            slide.shapes.add_picture(image_paths[i], Inches(3.5), Inches(4.5), width=Inches(2), height=Inches(2))
            content_textbox = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(16)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_split_full_image" and has_image:
            # Yarmi rasm, yarmi matn
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(1.5), width=Inches(5), height=Inches(5))
            content_textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_overlay_image_dark" and has_image:
            # Rasm ustiga qora shaffof qoplama, ustida matn
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.5
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(255, 255, 255)  # Oq matn

        elif style == "style_diagonal_layout" and has_image:
            # Diagonal bo‘linma
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(1.5), width=Inches(5), height=Inches(5))
            content_textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_focus_image_blend" and has_image:
            # Butun fon rasm, ustida shaffof matn bloki
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(4))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
            overlay.fill.transparency = 0.7
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_focus_circle" and has_image:
            # Markazda doira rasm, atrofida matn
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(2), Inches(3), Inches(3))
            circle.fill.background()
            slide.shapes.add_picture(image_paths[i], Inches(3.5), Inches(2), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.alignment = PP_ALIGN.CENTER

        elif style == "style_stat_block" and has_image:
            # Vizual statistika + rasm
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_side_by_side_card" and has_image:
            # Ikkita kartochka: biri rasm, biri matn
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(4), height=Inches(4))
            content_textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_background_with_quote" and has_image:
            # Rasm fonida iqtibos
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(4))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.5
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = f"“{slide_texts[i]}”"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        elif style == "style_basic_text_centered":
            # Markazlashtirilgan matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            for p in text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.alignment = PP_ALIGN.CENTER

        elif style == "style_bullet_points_left":
            # Chapga yo‘naltirilgan bullet points
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"• {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_numbered_steps_right":
            # O‘ngda raqamlangan bosqichlar
            content_textbox = slide.shapes.add_textbox(Inches(5), Inches(2.5), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"{idx}. {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_question_vs_answer":
            # Savol-javob shaklida bo‘linma
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = "Savol: " + slide_titles[i]
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.font.bold = True
            content_textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame2 = content_textbox2.text_frame
            text_frame2.word_wrap = True
            p2 = text_frame2.add_paragraph()
            p2.text = "Javob: " + slide_texts[i]
            p2.font.size = Pt(20)
            p2.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_vertical_steps":
            # Yuqoridan pastga bosqichlar
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"Bosqich {idx}: {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_quote_spotlight":
            # Spotlight effektida iqtibos
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = f"“{slide_texts[i]}”"
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True

        elif style == "style_simple_title_content":
            # Yuqorida sarlavha, pastda matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_two_column_compare":
            # Ikkita ustun: taqqoslash
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences[:len(sentences)//2]:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
            content_textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame2 = content_textbox2.text_frame
            text_frame2.word_wrap = True
            for sentence in sentences[len(sentences)//2:]:
                if sentence.strip():
                    p = text_frame2.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_centered_fact":
            # Bitta muhim fikr markazda
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = slide_texts[i]
            p.font.size = Pt(28)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True

        elif style == "style_timeline_wave":
            # Vaqt chizig‘i to‘lqin uslubida
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"{idx}-yil: {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_call_to_action_box":
            # Harakatga chaqiruv
            content_textbox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = slide_texts[i]
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True

        elif style == "style_title_with_subtext":
            # Katta sarlavha, kichik izoh
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = slide_titles[i]
            p.font.size = Pt(28)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p2 = text_frame.add_paragraph()
            p2.text = slide_texts[i]
            p2.font.size = Pt(18)
            p2.font.color.rgb = hex_to_rgb(text_color)
            p2.alignment = PP_ALIGN.CENTER

        elif style == "style_highlight_word":
            # Muhim so‘zlar rang bilan ajratilgan
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_process_chain":
            # Zanjir uslubida bosqichlar
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"-> {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_stack_blocks":
            # Vertikal stack-bloklar
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_minimalist_boxed":
            # Minimalist quti ichida matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_definitions":
            # Term + Ta’rif
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"Ta'rif {idx+1}: {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        else:
            # Default holat: faqat matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

    return ppt