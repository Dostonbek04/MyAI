from pptx import Presentation as PPT
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math
import random

# Dizayn uslublari ro‘yxatlari (ma'lumot uchun)
IMAGE_STYLES = [
    "style_focus_left_image", "style_focus_right_image", "style_image_top_text_bottom",
    "style_image_grid_blocks", "style_split_full_image", "style_overlay_image_dark",
    "style_diagonal_layout", "style_focus_image_blend", "style_focus_circle",
    "style_stat_block", "style_side_by_side_card", "style_background_with_quote"
]
TEXT_STYLES = [
    "style_basic_text_centered", "style_bullet_points_left", "style_numbered_steps_right",
    "style_question_vs_answer", "style_vertical_steps", "style_quote_spotlight",
    "style_simple_title_content", "style_two_column_compare", "style_centered_fact",
    "style_timeline_wave", "style_call_to_action_box", "style_title_with_subtext",
    "style_highlight_word", "style_process_chain", "style_stack_blocks",
    "style_minimalist_boxed", "style_definitions", "style_circle_with_texts"
]

# Barcha style’lar ro‘yxati (noyob tanlash uchun)
ALL_STYLES = IMAGE_STYLES + TEXT_STYLES

# Dark shablonlar ro‘yxati
DARK_TEMPLATES = [
    {"name": "Tungi Osmon", "background_color": "#1c2526", "text_color": "#ffffff", "style": "style_focus_left_image"},
    {"name": "Qorong‘u Ufq", "background_color": "#2e2e2e", "text_color": "#d3d3d3", "style": "style_image_top_text_bottom"},
    {"name": "Sirli Tong", "background_color": "#1a1a1a", "text_color": "#e0e0e0", "style": "style_basic_text_centered"},
    {"name": "Qorong‘u Shamol", "background_color": "#252525", "text_color": "#cccccc", "style": "style_split_full_image"},
    {"name": "Tungi Dunyo", "background_color": "#121212", "text_color": "#ffffff", "style": "style_bullet_points_left"},
    {"name": "Qora Ufq", "background_color": "#1f1f1f", "text_color": "#d9d9d9", "style": "style_numbered_steps_right"},
    {"name": "Sirli Kecha", "background_color": "#2a2a2a", "text_color": "#e5e5e5", "style": "style_question_vs_answer"},
    {"name": "Tungi Yulduzlar", "background_color": "#151515", "text_color": "#f0f0f0", "style": "style_vertical_steps"},
    {"name": "Qorong‘u Yomg‘ir", "background_color": "#222222", "text_color": "#d0d0d0", "style": "style_quote_spotlight"},
    {"name": "Tungi Chiroq", "background_color": "#1d1d1d", "text_color": "#e8e8e8", "style": "style_simple_title_content"},
    {"name": "Qora Quyosh", "background_color": "#282828", "text_color": "#d5d5d5", "style": "style_two_column_compare"},
    {"name": "Tungi Tuman", "background_color": "#171717", "text_color": "#f5f5f5", "style": "style_centered_fact"},
    {"name": "Qorong‘u O‘rmon", "background_color": "#202020", "text_color": "#dadada", "style": "style_timeline_wave"},
    {"name": "Tungi Tong", "background_color": "#262626", "text_color": "#e2e2e2", "style": "style_call_to_action_box"},
    {"name": "Qora Suvlar", "background_color": "#141414", "text_color": "#f8f8f8", "style": "style_title_with_subtext"},
    {"name": "Tungi Nur", "background_color": "#232323", "text_color": "#d8d8d8", "style": "style_highlight_word"},
    {"name": "Qorong‘u Xazina", "background_color": "#1e1e1e", "text_color": "#e6e6e6", "style": "style_process_chain"},
    {"name": "Qora Shahar", "background_color": "#292929", "text_color": "#d4d4d4", "style": "style_stack_blocks"},
    {"name": "Tungi O‘tlar", "background_color": "#161616", "text_color": "#f2f2f2", "style": "style_minimalist_boxed"},
    {"name": "Qorong‘u Soyalar", "background_color": "#212121", "text_color": "#dedede", "style": "style_definitions"},
    {"name": "Qora Osmon", "background_color": "#131313", "text_color": "#f6f6f6", "style": "style_focus_right_image"},
    {"name": "Tungi Gullar", "background_color": "#242424", "text_color": "#d6d6d6", "style": "style_image_grid_blocks"},
    {"name": "Qorong‘u Qanotlar", "background_color": "#1b1b1b", "text_color": "#e4e4e4", "style": "style_overlay_image_dark"},
    {"name": "Qora Ufq 2", "background_color": "#272727", "text_color": "#e1e1e1", "style": "style_diagonal_layout"},
    {"name": "Tungi Chiroq 2", "background_color": "#181818", "text_color": "#f4f4f4", "style": "style_focus_image_blend"},
    {"name": "Qorong‘u Tushlar", "background_color": "#2c2c2c", "text_color": "#d2d2d2", "style": "style_focus_circle"},
    {"name": "Qora Bahor", "background_color": "#191919", "text_color": "#e9e9e9", "style": "style_stat_block"},
    {"name": "Tungi Soyalar", "background_color": "#252525", "text_color": "#dddddd", "style": "style_side_by_side_card"},
    {"name": "Qora Yulduzlar", "background_color": "#1a1a1a", "text_color": "#f1f1f1", "style": "style_background_with_quote"},
    {"name": "Tungi Doira", "background_color": "#1e1e1e", "text_color": "#ffffff", "style": "style_circle_with_texts"},
]

def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def set_background_color(slide, hex_color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(hex_color)

def create_dark_template(title, slide_titles, slide_texts, image_paths, num_slides, style_index, with_images=False, slide_styles=None):
    if style_index < 0 or style_index >= len(DARK_TEMPLATES):
        style_index = 0
    template = DARK_TEMPLATES[style_index]
    background_color = template["background_color"]
    text_color = template["text_color"]

    # Agar slide_styles berilmagan bo‘lsa, har bir slayd uchun noyob style tanlanadi
    if slide_styles is None or len(slide_styles) != num_slides:
        used_styles = []  # Tanlangan style’larni saqlash uchun ro‘yxat
        slide_styles = []
        available_styles = ALL_STYLES.copy()  # Barcha style’lardan nusxa olamiz

        for _ in range(num_slides):
            if not available_styles:  # Agar style’lar tugasa, ro‘yxatni qayta boshlaymiz
                available_styles = ALL_STYLES.copy()
                used_styles = []
            # Tasodifiy noyob style tanlash
            style = random.choice(available_styles)
            available_styles.remove(style)  # Tanlangan style’ni ro‘yxatdan o‘chiramiz
            used_styles.append(style)  # Tanlangan style’ni saqlaymiz
            slide_styles.append(style)

    ppt = PPT()

    # Birinchi slayd (sarlavha slaydi)
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    set_background_color(slide, background_color)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.left = Inches(1)
    title_shape.top = Inches(1)
    title_shape.width = Inches(8)
    title_shape.height = Inches(1)
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(text_color)

    # Qolgan slaydlar
    for i in range(num_slides):
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        set_background_color(slide, background_color)

        # Sarlavha har doim yuqorida joylashadi
        title_shape = slide.shapes.title
        title_shape.text = slide_titles[i]
        title_shape.left = Inches(1)
        title_shape.top = Inches(1)
        title_shape.width = Inches(8)
        title_shape.height = Inches(1)
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = hex_to_rgb(text_color)

        has_image = with_images and image_paths[i] is not None
        style = slide_styles[i]

        # Dizayn uslubiga ko‘ra slaydni joylashtirish
        if style == "style_focus_left_image" and has_image:
            # Chapda rasm, o‘ngda matn
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2.5), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(5), Inches(2.5), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_focus_right_image" and has_image:
            # O‘ngda rasm, chapda matn
            slide.shapes.add_picture(image_paths[i], Inches(6), Inches(2.5), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_image_top_text_bottom" and has_image:
            # Yuqorida rasm, pastda matn
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(8), height=Inches(2.5))
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_image_grid_blocks" and has_image:
            # 2x2 rasmli bloklar
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(2), height=Inches(2))
            slide.shapes.add_picture(image_paths[i], Inches(3.5), Inches(2), width=Inches(2), height=Inches(2))
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(4.5), width=Inches(2), height=Inches(2))
            slide.shapes.add_picture(image_paths[i], Inches(3.5), Inches(4.5), width=Inches(2), height=Inches(2))
            content_textbox = slide.shapes.add_textbox(Inches(6), Inches(2), Inches(3), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(16)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_split_full_image" and has_image:
            # Yarmi rasm, yarmi matn
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(1.5), width=Inches(5), height=Inches(5))
            content_textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_overlay_image_dark" and has_image:
            # Rasm ustiga qora shaffof qoplama, ustida matn
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.5
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(255, 255, 255)  # Oq matn

        elif style == "style_diagonal_layout" and has_image:
            # Diagonal bo‘linma
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(1.5), width=Inches(5), height=Inches(5))
            content_textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_focus_image_blend" and has_image:
            # Butun fon rasm, ustida shaffof matn bloki
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(4))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
            overlay.fill.transparency = 0.7
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_focus_circle" and has_image:
            # Markazda doira rasm, atrofida matn
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(2), Inches(3), Inches(3))
            circle.fill.background()
            slide.shapes.add_picture(image_paths[i], Inches(3.5), Inches(2), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.alignment = PP_ALIGN.CENTER

        elif style == "style_stat_block" and has_image:
            # Vizual statistika + rasm
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(3), height=Inches(3))
            content_textbox = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_side_by_side_card" and has_image:
            # Ikkita kartochka: biri rasm, biri matn
            slide.shapes.add_picture(image_paths[i], Inches(1), Inches(2), width=Inches(4), height=Inches(4))
            content_textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_background_with_quote" and has_image:
            # Rasm fonida iqtibos
            slide.shapes.add_picture(image_paths[i], Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(4))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.5
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = f"“{slide_texts[i]}”"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        elif style == "style_circle_with_texts":
            # Markazda doira ichida bitta so‘z, atrofida kichik matnlar
            # 1. Markazdagi doira va matn
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(2), Inches(3), Inches(3))
            circle.fill.solid()
            circle.fill.fore_color.rgb = hex_to_rgb("#333333")  # Qorong‘i kulrang doira foni
            content_textbox = slide.shapes.add_textbox(Inches(3.5), Inches(2), Inches(3), Inches(3))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            # Faqat birinchi so‘zni olish
            first_word = slide_texts[i].split()[0] if slide_texts[i].split() else "Idea"
            p.text = first_word
            p.font.size = Pt(28)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            # Vertikal markazlashtirish uchun
            text_frame.paragraphs[0].space_before = Pt(60)

            # 2. Atrafdagi kichik matnlar
            # Matnlar sonini aniqlash (3, 4 yoki 5)
            sentences = slide_texts[i].split('. ')
            short_texts = []
            for sentence in sentences:
                words = sentence.split()
                # Har bir matn 1-2 so‘zdan iborat bo‘ladi
                if words:
                    short_texts.append(' '.join(words[:2]))
            num_texts = min(len(short_texts), 5)  # Maksimum 5 ta matn
            if num_texts < 3:
                num_texts = 3  # Minimal 3 ta matn bo‘lishi kerak
                while len(short_texts) < num_texts:
                    short_texts.append("Qisqa Matn")  # Agar matn yetishmasa, standart qo‘shiladi

            # Burchaklarni hisoblash
            radius = 2  # Doira radiusi (inch)
            center_x, center_y = 5, 3.5  # Markaziy nuqta (slaydning o‘rtasi)
            text_box_width, text_box_height = 1.5, 0.5  # Kichik matn qutisi o‘lchami (inch)
            for idx in range(num_texts):
                angle = (360 / num_texts) * idx  # Har bir matnning burchagi
                rad = math.radians(angle)
                # Matn qutisining pozitsiyasini hisoblash
                text_x = center_x + radius * math.cos(rad) - (text_box_width / 2)
                text_y = center_y + radius * math.sin(rad) - (text_box_height / 2)
                # Kichik matn qutisini qo‘shish
                content_textbox = slide.shapes.add_textbox(
                    Inches(text_x), Inches(text_y), Inches(text_box_width), Inches(text_box_height)
                )
                text_frame = content_textbox.text_frame
                text_frame.word_wrap = True
                p = text_frame.add_paragraph()
                p.text = short_texts[idx]
                p.font.size = Pt(14)  # Kichikroq shrift
                p.font.color.rgb = hex_to_rgb(text_color)
                p.alignment = PP_ALIGN.CENTER

        elif style == "style_basic_text_centered":
            # Markazlashtirilgan matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            for p in text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.alignment = PP_ALIGN.CENTER

        elif style == "style_bullet_points_left":
            # Chapga yo‘naltirilgan bullet points
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"• {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_numbered_steps_right":
            # O‘ngda raqamlangan bosqichlar
            content_textbox = slide.shapes.add_textbox(Inches(5), Inches(2.5), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"{idx}. {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_question_vs_answer":
            # Savol-javob shaklida bo‘linma
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4.5), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = "Savol: " + slide_titles[i]
            p.font.size = Pt(20)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.font.bold = True
            content_textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame2 = content_textbox2.text_frame
            text_frame2.word_wrap = True
            p2 = text_frame2.add_paragraph()
            p2.text = "Javob: " + slide_texts[i]
            p2.font.size = Pt(20)
            p2.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_vertical_steps":
            # Yuqoridan pastga bosqichlar
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"Bosqich {idx}: {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_quote_spotlight":
            # Spotlight effektida iqtibos
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = f"“{slide_texts[i]}”"
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True

        elif style == "style_simple_title_content":
            # Yuqorida sarlavha, pastda matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_two_column_compare":
            # Ikkita ustun: taqqoslash
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences[:len(sentences)//2]:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
            content_textbox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
            text_frame2 = content_textbox2.text_frame
            text_frame2.word_wrap = True
            for sentence in sentences[len(sentences)//2:]:
                if sentence.strip():
                    p = text_frame2.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_centered_fact":
            # Bitta muhim fikr markazda
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = slide_texts[i]
            p.font.size = Pt(28)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True

        elif style == "style_timeline_wave":
            # Vaqt chizig‘i to‘lqin uslubida
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"{idx}-yil: {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_call_to_action_box":
            # Harakatga chaqiruv
            content_textbox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(2))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = slide_texts[i]
            p.font.size = Pt(24)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True

        elif style == "style_title_with_subtext":
            # Katta sarlavha, kichik izoh
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = slide_titles[i]
            p.font.size = Pt(28)
            p.font.color.rgb = hex_to_rgb(text_color)
            p.alignment = PP_ALIGN.CENTER
            p.font.bold = True
            p2 = text_frame.add_paragraph()
            p2.text = slide_texts[i]
            p2.font.size = Pt(18)
            p2.font.color.rgb = hex_to_rgb(text_color)
            p2.alignment = PP_ALIGN.CENTER

        elif style == "style_highlight_word":
            # Muhim so‘zlar rang bilan ajratilgan
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

        elif style == "style_process_chain":
            # Zanjir uslubida bosqichlar
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences, 1):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"-> {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_stack_blocks":
            # Vertikal stack-bloklar
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_minimalist_boxed":
            # Minimalist quti ichida matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        elif style == "style_definitions":
            # Term + Ta’rif
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for idx, sentence in enumerate(sentences):
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = f"Ta'rif {idx+1}: {sentence.strip()}" + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)
                    p.space_before = Pt(10)

        else:
            # Default holat: faqat matn
            content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(4))
            text_frame = content_textbox.text_frame
            text_frame.word_wrap = True
            sentences = slide_texts[i].split('. ')
            for sentence in sentences:
                if sentence.strip():
                    p = text_frame.add_paragraph()
                    p.text = sentence.strip() + ('.' if not sentence.endswith('.') else '')
                    p.font.size = Pt(20)
                    p.font.color.rgb = hex_to_rgb(text_color)

    return ppt