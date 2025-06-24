import openai
import requests
import os
import subprocess
import tempfile
from pptx import Presentation
from .presentationsTemplates.templates.dark_templates import create_dark_template
from .presentationsTemplates.templates.light_templates import create_light_template
from .presentationsTemplates.templates.professional_templates import create_professional_template
from PIL import Image, ImageDraw, ImageFont
from django.conf import settings
from pptx.util import Inches
from pdf2image import convert_from_path  # Yangilangan import
import dropbox
from dropbox.exceptions import ApiError
from django.conf import settings
from PIL import Image, ImageDraw, ImageFont
import os

# OpenAI API kalitini sozlash
openai.api_key = "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"

from django.core.files.base import ContentFile

def generate_slide_text_with_openai(title, num_slides=8):
    """
    OpenAI yordamida barcha slaydlar uchun matn generatsiya qilish.
    Har bir slaydda 10 ta list (sarlavha + qisqa tavsif) bo‘ladi.
    """
    prompt = (
        f"{title} mavzusida {num_slides} ta slayd uchun matnlar generatsiya qil. "
        f"Har bir slayd uchun 10 ta list (har biri sarlavha va 2-3 jumlali qisqa tavsifdan iborat) yoz. "
        f"Har bir list yangi qatordan boshlansin, slaydlar esa bo‘sh qator bilan ajratilsin."
    )
    try:
        client = openai.OpenAI(api_key=openai.api_key)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are a helpful assistant for generating presentation content."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1000,
            temperature=0.7
        )
        slide_texts = response.choices[0].message.content.split("\n\n")
        while len(slide_texts) < num_slides:
            slide_texts.append("Tavsif yo‘q")
        return slide_texts[:num_slides]
    except Exception as e:
        print(f"Matn generatsiya qilishda xato: {e}")
        return ["Tavsif yo‘q"] * num_slides

def generate_images_with_openai(title, num_slides, slide_texts):
    """
    OpenAI DALL·E yordamida har bir slayd uchun 5 ta rasm generatsiya qilish.
    """
    image_paths = []
    for i in range(num_slides):
        for j in range(5):
            slide_lines = slide_texts[i].split("\n")
            keywords = slide_lines[j*2] if j*2 < len(slide_lines) else "Umumiy himoya vositasi"
            prompt = f"{title} mavzusida slayd uchun rasm: {keywords}"

            try:
                client = openai.OpenAI(api_key=openai.api_key)
                response = client.images.generate(
                    prompt=prompt,
                    n=1,
                    size="1024x1024",
                    quality="standard"
                )
                image_url = response.data[0].url
                image_response = requests.get(image_url)
                image_response.raise_for_status()
                image_path = f"media/presentation_images/{title}_slide_{i+1}_image_{j+1}.jpg"
                os.makedirs(os.path.dirname(image_path), exist_ok=True)
                with open(image_path, 'wb') as f:
                    f.write(image_response.content)
                image_paths.append(image_path)
            except Exception as e:
                print(f"Rasm generatsiyasida xato (slayd {i+1}, rasm {j+1}): {e}")
                image_paths.append(None)

    return image_paths

def create_template_preview(template_name, background_color, text_color, template_type):
    """
    Shablon uchun preview rasm yaratish.
    """
    try:
        # Rasm o'lchami
        width, height = 800, 600
        image = Image.new('RGB', (width, height), background_color)
        draw = ImageDraw.Draw(image)

        # Matnlar qo'shish
        try:
            font_large = ImageFont.truetype("arial.ttf", 40)
            font_small = ImageFont.truetype("arial.ttf", 20)
        except:
            font_large = ImageFont.load_default()
            font_small = ImageFont.load_default()

        # Title matni
        draw.text((50, 50), "Title", fill=text_color, font=font_large)
        # Body & Link matni
        draw.text((50, 150), "Body & Link", fill=text_color, font=font_small)

        # Rasmni saqlash
        preview_dir = os.path.join("media", "templates", "previews")
        os.makedirs(preview_dir, exist_ok=True)
        template_name_clean = template_name.replace("‘", "_").replace(" ", "_").lower()
        output_path = os.path.join(preview_dir, f"{template_type}_{template_name_clean}.jpg")
        image.save(output_path, "JPEG")
        return output_path
    except Exception as e:
        print(f"Rasm yaratishda xato: {template_type}_{template_name_clean} - Xato: {str(e)}")
        raise

class PPT:
    def __init__(self):
        self.prs = Presentation()

    def add_slide(self, title, content=None, image_path=None):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        if content:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content

        if image_path:
            slide.shapes.add_picture(image_path, Inches(1), Inches(2), width=Inches(3))

    def save(self, file_path):
        self.prs.save(file_path)

    def save(self, output_stream):
        self.prs.save(output_stream)

def generate_slide_preview(pptx_path, output_path):
    """
    .pptx faylning birinchi slaydini JPEG rasm sifatida saqlaydi.
    Args:
        pptx_path (str): .pptx fayl yo‘li
        output_path (str): Saqlanadigan rasm yo‘li (masalan, media/previews/1_preview.jpg)
    Returns:
        bool: Muvaffaqiyatli bo‘lsa True, aks holda False
    """
    try:
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
            temp_pdf_path = temp_pdf.name

        subprocess.run(
            ["unoconv", "-f", "pdf", "-o", temp_pdf_path, str(pptx_path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )

        images = convert_from_path(temp_pdf_path, first_page=1, last_page=1, dpi=200)
        if images:
            images[0].save(output_path, "JPEG", quality=80)

        os.unlink(temp_pdf_path)
        return True
    except (subprocess.CalledProcessError, Exception) as e:
        print(f"Rasm generatsiyasida xato: {str(e)}")
        return False

def upload_to_dropbox(file_path, dropbox_path):
    """Faylni Dropbox’ga yuklash"""
    dbx = dropbox.Dropbox(settings.DROPBOX_ACCESS_TOKEN)
    try:
        with open(file_path, 'rb') as f:
            dbx.files_upload(f.read(), dropbox_path, mode=dropbox.files.WriteMode('overwrite'))
        # Dropbox’dan umumiy link olish
        shared_link_metadata = dbx.sharing_create_shared_link_with_settings(dropbox_path)
        return shared_link_metadata.url
    except ApiError as e:
        print(f"Dropbox yuklashda xato: {e}")
        return None

def delete_from_dropbox(dropbox_path):
    """Dropbox’dan faylni o‘chirish"""
    dbx = dropbox.Dropbox(settings.DROPBOX_ACCESS_TOKEN)
    try:
        dbx.files_delete_v2(dropbox_path)
    except ApiError as e:
        print(f"Dropbox o‘chirishda xato: {e}")